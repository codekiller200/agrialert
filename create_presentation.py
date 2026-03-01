"""
Script de génération de la présentation PRESCI pour AgriAlert BF
Crée une présentation PowerPoint professionnelle avec un design inspiré du Burkina Faso
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_agrialert_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Palette de couleurs AgriAlert BF (inspirée du Burkina Faso)
    PRIMARY_GREEN = RGBColor(45, 80, 22)      # #2D5016
    ACCENT_ORANGE = RGBColor(224, 123, 57)    # #E07B39
    GOLD_YELLOW = RGBColor(244, 196, 48)      # #F4C430
    WARM_BEIGE = RGBColor(245, 230, 211)      # #F5E6D3
    DEEP_BROWN = RGBColor(92, 64, 51)         # #5C4033
    WHITE = RGBColor(255, 255, 255)
    
    # SLIDE 1: Titre principal
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Fond dégradé vert
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_GREEN
    
    # Titre principal
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "AgriAlert BF"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = 'Georgia'
    
    # Sous-titre
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(1.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Application mobile d'alerte sécheresse\npour les agriculteurs burkinabè"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = GOLD_YELLOW
    subtitle_para.font.name = 'Calibri'
    
    # Badge PRESCI
    presci_box = slide1.shapes.add_textbox(
        Inches(3.5), Inches(6), Inches(3), Inches(0.6)
    )
    presci_frame = presci_box.text_frame
    presci_frame.text = "PRESCI 2026 - Niveau Universitaire"
    presci_para = presci_frame.paragraphs[0]
    presci_para.alignment = PP_ALIGN.CENTER
    presci_para.font.size = Pt(14)
    presci_para.font.color.rgb = WHITE
    presci_para.font.name = 'Calibri'
    
    # SLIDE 2: Problématique
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = WARM_BEIGE
    
    # Bandeau titre
    title_shape2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape2.fill.solid()
    title_shape2.fill.fore_color.rgb = PRIMARY_GREEN
    title_shape2.line.fill.background()
    
    # Titre
    title_text2 = title_shape2.text_frame
    title_text2.text = "La Problématique"
    title_text2.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_text2.paragraphs[0].font.size = Pt(40)
    title_text2.paragraphs[0].font.bold = True
    title_text2.paragraphs[0].font.color.rgb = WHITE
    title_text2.paragraphs[0].font.name = 'Georgia'
    title_text2.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Contenu - 3 cartes problèmes
    problems = [
        {
            "title": "Sécheresses Récurrentes",
            "text": "Le Burkina Faso fait face à des sécheresses de plus en plus fréquentes, menaçant la sécurité alimentaire de millions de personnes."
        },
        {
            "title": "Manque d'Information",
            "text": "Les agriculteurs ruraux n'ont pas accès à des prévisions météorologiques fiables et des alertes précoces adaptées à leur contexte."
        },
        {
            "title": "Connectivité Limitée",
            "text": "Les zones agricoles souffrent de connexions internet instables, limitant l'accès aux services numériques d'aide à la décision."
        }
    ]
    
    y_start = 1.8
    for i, problem in enumerate(problems):
        # Carte
        card = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_start), Inches(8.4), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ACCENT_ORANGE
        card.line.width = Pt(2)
        
        # Numéro
        num_shape = slide2.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), Inches(y_start + 0.3), Inches(0.7), Inches(0.7)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = ACCENT_ORANGE
        num_shape.line.fill.background()
        num_text = num_shape.text_frame
        num_text.text = str(i + 1)
        num_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_text.paragraphs[0].font.size = Pt(28)
        num_text.paragraphs[0].font.bold = True
        num_text.paragraphs[0].font.color.rgb = WHITE
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Titre du problème
        problem_title = slide2.shapes.add_textbox(
            Inches(2.2), Inches(y_start + 0.1), Inches(6.8), Inches(0.5)
        )
        pf = problem_title.text_frame
        pf.text = problem["title"]
        pf.paragraphs[0].font.size = Pt(20)
        pf.paragraphs[0].font.bold = True
        pf.paragraphs[0].font.color.rgb = DEEP_BROWN
        pf.paragraphs[0].font.name = 'Calibri'
        
        # Texte du problème
        problem_text = slide2.shapes.add_textbox(
            Inches(2.2), Inches(y_start + 0.6), Inches(6.8), Inches(0.6)
        )
        ptf = problem_text.text_frame
        ptf.text = problem["text"]
        ptf.paragraphs[0].font.size = Pt(14)
        ptf.paragraphs[0].font.color.rgb = DEEP_BROWN
        ptf.paragraphs[0].font.name = 'Calibri'
        ptf.word_wrap = True
        
        y_start += 1.6
    
    # SLIDE 3: Notre Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    background3 = slide3.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = WARM_BEIGE
    
    # Bandeau titre
    title_shape3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape3.fill.solid()
    title_shape3.fill.fore_color.rgb = ACCENT_ORANGE
    title_shape3.line.fill.background()
    
    title_text3 = title_shape3.text_frame
    title_text3.text = "Notre Solution: AgriAlert BF"
    title_text3.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_text3.paragraphs[0].font.size = Pt(40)
    title_text3.paragraphs[0].font.bold = True
    title_text3.paragraphs[0].font.color.rgb = WHITE
    title_text3.paragraphs[0].font.name = 'Georgia'
    title_text3.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Description principale
    desc_box = slide3.shapes.add_textbox(
        Inches(1), Inches(1.8), Inches(8), Inches(1)
    )
    desc_frame = desc_box.text_frame
    desc_frame.text = "Une application mobile utilisant l'intelligence artificielle pour prédire les risques de sécheresse et fournir des conseils agricoles personnalisés, fonctionnant même sans connexion internet."
    desc_para = desc_frame.paragraphs[0]
    desc_para.alignment = PP_ALIGN.CENTER
    desc_para.font.size = Pt(18)
    desc_para.font.color.rgb = DEEP_BROWN
    desc_para.font.name = 'Calibri'
    desc_frame.word_wrap = True
    
    # Fonctionnalités clés - grille 2x2
    features = [
        ("🤖 IA Embarquée", "Modèle TensorFlow Lite\npour prédictions offline"),
        ("📍 Géolocalisation", "GPS automatique ou\nsélection manuelle"),
        ("🌦️ Météo 7 jours", "Prévisions détaillées\nAPI Open-Meteo"),
        ("💡 Conseils Pratiques", "Recommandations\nadaptées au risque")
    ]
    
    x_positions = [1.2, 5.4]
    y_positions = [3.2, 5.2]
    
    for i, feature in enumerate(features):
        row = i // 2
        col = i % 2
        x = x_positions[col]
        y = y_positions[row]
        
        # Carte fonctionnalité
        feat_card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(3.6), Inches(1.6)
        )
        feat_card.fill.solid()
        feat_card.fill.fore_color.rgb = WHITE
        feat_card.line.color.rgb = PRIMARY_GREEN
        feat_card.line.width = Pt(2)
        
        # Titre fonctionnalité
        feat_title = slide3.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.2), Inches(3.2), Inches(0.4)
        )
        ftf = feat_title.text_frame
        ftf.text = feature[0]
        ftf.paragraphs[0].font.size = Pt(18)
        ftf.paragraphs[0].font.bold = True
        ftf.paragraphs[0].font.color.rgb = PRIMARY_GREEN
        ftf.paragraphs[0].font.name = 'Calibri'
        ftf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Description fonctionnalité
        feat_desc = slide3.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.7), Inches(3.2), Inches(0.8)
        )
        fdf = feat_desc.text_frame
        fdf.text = feature[1]
        fdf.paragraphs[0].font.size = Pt(14)
        fdf.paragraphs[0].font.color.rgb = DEEP_BROWN
        fdf.paragraphs[0].font.name = 'Calibri'
        fdf.paragraphs[0].alignment = PP_ALIGN.CENTER
        fdf.word_wrap = True
    
    # SLIDE 4: Technologies Utilisées
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background4 = slide4.background
    fill4 = background4.fill
    fill4.solid()
    fill4.fore_color.rgb = WARM_BEIGE
    
    title_shape4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_shape4.fill.solid()
    title_shape4.fill.fore_color.rgb = PRIMARY_GREEN
    title_shape4.line.fill.background()
    
    title_text4 = title_shape4.text_frame
    title_text4.text = "Technologies Utilisées"
    title_text4.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_text4.paragraphs[0].font.size = Pt(40)
    title_text4.paragraphs[0].font.bold = True
    title_text4.paragraphs[0].font.color.rgb = WHITE
    title_text4.paragraphs[0].font.name = 'Georgia'
    title_text4.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Technologies - 3 colonnes
    techs = [
        {
            "category": "Frontend",
            "items": ["Flutter 3.0+", "Dart", "Provider", "Flutter Map"]
        },
        {
            "category": "IA & Backend",
            "items": ["TensorFlow Lite", "Python", "scikit-learn", "Open-Meteo API"]
        },
        {
            "category": "Données & Services",
            "items": ["OpenStreetMap", "Geolocator", "Shared Preferences"]
        }
    ]
    
    x_start = 0.8
    col_width = 2.8
    for i, tech in enumerate(techs):
        x = x_start + (i * 3.1)
        
        # Carte catégorie
        cat_card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2), Inches(col_width), Inches(4)
        )
        cat_card.fill.solid()
        cat_card.fill.fore_color.rgb = WHITE
        cat_card.line.color.rgb = ACCENT_ORANGE
        cat_card.line.width = Pt(2)
        
        # Titre catégorie
        cat_title = slide4.shapes.add_textbox(
            Inches(x + 0.2), Inches(2.2), Inches(col_width - 0.4), Inches(0.5)
        )
        ctf = cat_title.text_frame
        ctf.text = tech["category"]
        ctf.paragraphs[0].font.size = Pt(20)
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].font.color.rgb = ACCENT_ORANGE
        ctf.paragraphs[0].font.name = 'Calibri'
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Items
        y_item = 3
        for item in tech["items"]:
            item_box = slide4.shapes.add_textbox(
                Inches(x + 0.3), Inches(y_item), Inches(col_width - 0.6), Inches(0.4)
            )
            itf = item_box.text_frame
            itf.text = "• " + item
            itf.paragraphs[0].font.size = Pt(14)
            itf.paragraphs[0].font.color.rgb = DEEP_BROWN
            itf.paragraphs[0].font.name = 'Calibri'
            y_item += 0.6
    
    # SLIDE 5: Impact et Innovation
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background5 = slide5.background
    fill5 = background5.fill
    fill5.solid()
    fill5.fore_color.rgb = PRIMARY_GREEN
    
    title_box5 = slide5.shapes.add_textbox(
        Inches(1), Inches(0.8), Inches(8), Inches(0.8)
    )
    tf5 = title_box5.text_frame
    tf5.text = "Impact et Innovation"
    tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf5.paragraphs[0].font.size = Pt(44)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = GOLD_YELLOW
    tf5.paragraphs[0].font.name = 'Georgia'
    
    # Stats d'impact
    stats = [
        ("2M+", "Agriculteurs potentiellement aidés"),
        ("100%", "Fonctionnement offline"),
        ("2 langues", "Français & Mooré")
    ]
    
    y_stat = 2.2
    for stat in stats:
        # Nombre
        num_box = slide5.shapes.add_textbox(
            Inches(1.5), Inches(y_stat), Inches(7), Inches(0.8)
        )
        nf = num_box.text_frame
        nf.text = stat[0]
        nf.paragraphs[0].alignment = PP_ALIGN.CENTER
        nf.paragraphs[0].font.size = Pt(60)
        nf.paragraphs[0].font.bold = True
        nf.paragraphs[0].font.color.rgb = GOLD_YELLOW
        nf.paragraphs[0].font.name = 'Impact'
        
        # Label
        label_box = slide5.shapes.add_textbox(
            Inches(1.5), Inches(y_stat + 0.7), Inches(7), Inches(0.4)
        )
        lf = label_box.text_frame
        lf.text = stat[1]
        lf.paragraphs[0].alignment = PP_ALIGN.CENTER
        lf.paragraphs[0].font.size = Pt(18)
        lf.paragraphs[0].font.color.rgb = WHITE
        lf.paragraphs[0].font.name = 'Calibri'
        
        y_stat += 1.5
    
    # SLIDE 6: Conclusion
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background6 = slide6.background
    fill6 = background6.fill
    fill6.solid()
    fill6.fore_color.rgb = PRIMARY_GREEN
    
    # Message principal
    main_msg = slide6.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    mmf = main_msg.text_frame
    mmf.text = "AgriAlert BF\nAnticipons ensemble la sécheresse\npour une agriculture résiliente"
    mmf.paragraphs[0].alignment = PP_ALIGN.CENTER
    mmf.paragraphs[0].font.size = Pt(36)
    mmf.paragraphs[0].font.bold = True
    mmf.paragraphs[0].font.color.rgb = WHITE
    mmf.paragraphs[0].font.name = 'Georgia'
    mmf.word_wrap = True
    
    # Icône
    icon_box = slide6.shapes.add_textbox(
        Inches(4), Inches(4), Inches(2), Inches(1)
    )
    icf = icon_box.text_frame
    icf.text = "🌾"
    icf.paragraphs[0].alignment = PP_ALIGN.CENTER
    icf.paragraphs[0].font.size = Pt(80)
    
    # Footer
    footer_box = slide6.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(0.8)
    )
    ff = footer_box.text_frame
    ff.text = "Merci pour votre attention\nPRESCI 2026 - Développement d'applications mobiles"
    ff.paragraphs[0].alignment = PP_ALIGN.CENTER
    ff.paragraphs[0].font.size = Pt(16)
    ff.paragraphs[0].font.color.rgb = GOLD_YELLOW
    ff.paragraphs[0].font.name = 'Calibri'
    
    # Sauvegarder
    output_path = '/mnt/user-data/outputs/AgriAlert_BF_Presentation_PRESCI_2026.pptx'
    prs.save(output_path)
    print(f"Présentation créée avec succès: {output_path}")
    return output_path

if __name__ == "__main__":
    create_agrialert_presentation()
